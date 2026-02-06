from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Sample Presentation'
slide.placeholders[1].text = 'Testing PPTX to Markdown conversion'

# Slide 2: Bullet points
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Features'
tf = slide.placeholders[1].text_frame
tf.text = 'Key features:'
p = tf.add_paragraph()
p.text = 'Bullet point 1'
p.level = 1
p = tf.add_paragraph()
p.text = 'Bullet point 2'
p.level = 1
p = tf.add_paragraph()
p.text = 'Sub-point'
p.level = 2

# Slide 3: Content with table
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Data Table'
rows = cols = 3
left = top = Inches(2)
width = Inches(6)
height = Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

table.cell(0, 0).text = 'Header 1'
table.cell(0, 1).text = 'Header 2'
table.cell(0, 2).text = 'Header 3'
table.cell(1, 0).text = 'Row 1 Col 1'
table.cell(1, 1).text = 'Row 1 Col 2'
table.cell(1, 2).text = 'Row 1 Col 3'
table.cell(2, 0).text = 'Row 2 Col 1'
table.cell(2, 1).text = 'Row 2 Col 2'
table.cell(2, 2).text = 'Row 2 Col 3'

# Slide 4: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Conclusion'
slide.placeholders[1].text = 'This presentation demonstrates PPTX conversion capabilities.'

prs.save('test.pptx')
